from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLORS = {
    "navy": RGBColor(18, 52, 86),
    "teal": RGBColor(33, 126, 140),
    "orange": RGBColor(225, 122, 95),
    "gold": RGBColor(230, 184, 92),
    "bg": RGBColor(245, 247, 250),
    "text": RGBColor(33, 37, 41),
    "muted": RGBColor(98, 108, 122),
    "line": RGBColor(208, 214, 221),
    "white": RGBColor(255, 255, 255),
    "soft_blue": RGBColor(232, 240, 248),
    "soft_teal": RGBColor(230, 245, 246),
    "soft_orange": RGBColor(251, 238, 234),
}

FONT_NAME = "Microsoft YaHei"


def set_run_font(run, size, bold=False, color=None, font_name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color or COLORS["text"]


def add_full_bg(slide, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_top_band(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.78)
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = COLORS["navy"]

    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.14), Inches(8.2), Inches(0.38))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    set_run_font(run, 24, bold=True, color=COLORS["white"])

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(8.7), Inches(0.18), Inches(4.0), Inches(0.28))
        sub_tf = sub_box.text_frame
        sub_tf.clear()
        sp = sub_tf.paragraphs[0]
        sp.alignment = PP_ALIGN.RIGHT
        srun = sp.add_run()
        srun.text = subtitle
        set_run_font(srun, 10, color=RGBColor(226, 231, 236))


def add_footer(slide, text="汇报PPT_260416"):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(7.1), Inches(12.2), Inches(0.02)
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["line"]

    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.13), Inches(3.2), Inches(0.2))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_font(run, 9, color=COLORS["muted"])


def add_title_slide(slide, title, subtitle, tag):
    add_full_bg(slide, COLORS["bg"])

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.85), Inches(5.3), Inches(5.45)
    )
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["navy"]

    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(1.1), Inches(0.22), Inches(4.95)
    )
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = COLORS["gold"]

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(1.05), Inches(3.35), Inches(0.62)
    )
    badge.line.color.rgb = COLORS["teal"]
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLORS["soft_teal"]
    badge.text_frame.clear()
    p = badge.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = tag
    set_run_font(r, 15, bold=True, color=COLORS["teal"])

    title_box = slide.shapes.add_textbox(Inches(1.35), Inches(1.35), Inches(4.3), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_run_font(r1, 28, bold=True, color=COLORS["white"])

    sub_box = slide.shapes.add_textbox(Inches(1.35), Inches(4.6), Inches(4.1), Inches(1.1))
    sub_tf = sub_box.text_frame
    sub_tf.word_wrap = True
    p2 = sub_tf.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    set_run_font(r2, 14, color=RGBColor(225, 231, 240))

    note = slide.shapes.add_textbox(Inches(6.7), Inches(2.15), Inches(5.3), Inches(2.8))
    note_tf = note.text_frame
    note_tf.word_wrap = True
    note_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    for idx, line in enumerate(
        [
            "基于 docs/0_avsStructure/ 现有 HTML 内容压缩整理",
            "适合 5 到 8 分钟汇报",
            "保留结构判断、映射关系与下一步建议",
        ]
    ):
        p = note_tf.paragraphs[0] if idx == 0 else note_tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = line
        set_run_font(r, 17 if idx == 0 else 15, bold=(idx == 0), color=COLORS["text"])

    add_footer(slide)


def add_bullet_box(slide, left, top, width, height, title, bullets, fill_color, title_color=None):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = COLORS["line"]
    card.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.14), width - Inches(0.35), Inches(0.35))
    title_tf = title_box.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_font(r, 17, bold=True, color=title_color or COLORS["navy"])

    body_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.clear()
    for idx, bullet in enumerate(bullets):
        p = body_tf.paragraphs[0] if idx == 0 else body_tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        p.bullet = True
        r = p.add_run()
        r.text = bullet
        set_run_font(r, 14, color=COLORS["text"])


def add_callout(slide, left, top, width, height, title, body, fill_color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.line.color.rgb = fill_color
    shape.line.width = Pt(1.5)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]

    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_run_font(r1, 16, bold=True, color=fill_color)

    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = body
    set_run_font(r2, 13, color=COLORS["text"])


def add_stage_row(slide, top, stages):
    left = Inches(0.75)
    box_w = Inches(2.65)
    gap = Inches(0.18)
    box_h = Inches(1.25)
    for idx, stage in enumerate(stages):
        x = left + idx * (box_w + gap)
        fill = stage["fill"]
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, box_w, box_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill

        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = stage["title"]
        set_run_font(r1, 16, bold=True, color=COLORS["white"])

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = stage["body"]
        set_run_font(r2, 12, color=COLORS["white"])

        if idx < len(stages) - 1:
            x1 = x + box_w
            y = top + Inches(0.62)
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, x1 + Inches(0.02), y, x1 + gap - Inches(0.02), y
            )
            line.line.color.rgb = COLORS["muted"]
            line.line.width = Pt(2)


def build_gsc_deck(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title_slide(
        slide,
        "GSC 多样本 .gsc\n二进制格式整理汇报",
        "基于 GSC多样本二进制格式汇报-260416.html\n聚焦文件结构、GT 语义与风险点",
        "简洁版 / 6 页",
    )

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "1. 核心结论", "多样本 .gsc 当前实现")
    add_bullet_box(
        slide,
        Inches(0.7),
        Inches(1.15),
        Inches(5.95),
        Inches(4.95),
        "这次整理明确了什么",
        [
            "`.gsc` 可拆成顶层头、Main Archive、lossless part2、SDSL tail 四段。",
            "GT 与 other fields 是两条不同存储链，不能按同一种字段流理解。",
            "当前 writer 实际输出 `GSCF V2` row-block 目录，reader 同时兼容 legacy 读取。",
            "文档已从“结构表”推进到“可追溯到源码的格式说明”。",
        ],
        COLORS["soft_blue"],
    )
    add_callout(
        slide,
        Inches(7.0),
        Inches(1.35),
        Inches(5.55),
        Inches(1.35),
        "适用范围",
        "只针对当前仓库的多样本 `gsc compress / decompress` on-disk 格式，不覆盖单样本 gVCF，也不等同于 AVS 外层容器。",
        COLORS["teal"],
    )
    add_callout(
        slide,
        Inches(7.0),
        Inches(2.95),
        Inches(5.55),
        Inches(1.35),
        "最重要的判断",
        "GT payload 不是文本 GT；它本质上是两条 bitplane 经过变换后的稀疏索引流，恢复时还要结合 zero/copy 位图。",
        COLORS["orange"],
    )
    add_callout(
        slide,
        Inches(7.0),
        Inches(4.55),
        Inches(5.55),
        Inches(1.15),
        "对外协议化前提",
        "若未来要把格式正式发布，需要补顶层 `magic/version` 并移除 `bool`、`size_t` 这类 ABI 敏感类型。",
        COLORS["navy"],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "2. 文件整体布局", "先看四段式容器")
    stages = [
        {"title": "顶层头", "body": "mode_type\nother_fields_offset\nsdsl_offset", "fill": COLORS["navy"]},
        {"title": "Main Archive", "body": "chunk 索引\nheader/sample\nfixed-fields chunks", "fill": COLORS["teal"]},
        {"title": "lossless part2", "body": "FILTER/INFO/FORMAT\n多流容器", "fill": COLORS["orange"]},
        {"title": "SDSL tail", "body": "4 个 rrr_vector\nzero/copy 位图", "fill": COLORS["gold"]},
    ]
    add_stage_row(slide, Inches(1.55), stages)
    add_bullet_box(
        slide,
        Inches(0.8),
        Inches(3.15),
        Inches(6.05),
        Inches(2.95),
        "Main Archive 里重点关注",
        [
            "chunk 定位信息与 `chunks_min_pos / where_chrom`。",
            "`ploidy / max_block_rows / max_block_cols / n_col_blocks` 等全局参数。",
            "压缩 header、样本名、fixed-fields chunk payload。",
        ],
        COLORS["soft_teal"],
        title_color=COLORS["teal"],
    )
    add_bullet_box(
        slide,
        Inches(6.95),
        Inches(3.15),
        Inches(5.55),
        Inches(2.95),
        "reader 侧真实读取顺序",
        [
            "先读顶层头部。",
            "先跳到 `sdsl_offset` 读取 4 个尾部位图。",
            "再回到头后位置读取 Main Archive。",
            "lossless 时把 `[other_fields_offset, sdsl_offset)` 当 part2 容器。",
        ],
        COLORS["soft_orange"],
        title_color=COLORS["orange"],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "3. Chunk 与 GT 真实语义", "当前 writer 实际输出 GSCF V2")
    add_bullet_box(
        slide,
        Inches(0.7),
        Inches(1.15),
        Inches(5.8),
        Inches(5.2),
        "Fixed-fields chunk",
        [
            "chunk 内带 `GSCF` magic 和 row-block 目录。",
            "V2 目录显式记录 `variant_count / first_pos / last_pos / gt_off / gt_size`。",
            "reader 先尝试按新目录读；不命中时回退到 legacy 平铺格式。",
            "这是按范围读取、按块解码的基础。",
        ],
        COLORS["soft_blue"],
    )
    add_bullet_box(
        slide,
        Inches(6.8),
        Inches(1.15),
        Inches(5.8),
        Inches(5.2),
        "GT payload",
        [
            "GT 不是 `0/0` 文本，而是两条 bitplane 的变换结果。",
            "`00/01/10/11 -> 0/1/./2`，因此 `no_vec` 近似 `2 × variants`。",
            "落盘前会经过置换、XOR、稀疏化与 copy/zero 标注。",
            "row-block GT payload 末尾还会追加 backend marker，reader 取出 segment 后再按 marker 解码。",
        ],
        COLORS["soft_teal"],
        title_color=COLORS["teal"],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "4. Lossless Part2 与查询路径", "other fields 是多流容器")
    add_bullet_box(
        slide,
        Inches(0.7),
        Inches(1.15),
        Inches(5.95),
        Inches(5.1),
        "Part2 的组织方式",
        [
            "`File_Handle_2` 管理多个 stream，而不是单一顺序字段流。",
            "每个 key 通常注册 `size` 与 `data` 两条流，另有 `part2_params` 存全局参数。",
            "`part2_params` 中还会放 `actual_variants`、`keys`、backend 与 FORMAT dictionaries。",
        ],
        COLORS["soft_orange"],
        title_color=COLORS["orange"],
    )
    add_bullet_box(
        slide,
        Inches(6.95),
        Inches(1.15),
        Inches(5.55),
        Inches(5.1),
        "查询与恢复时的价值",
        [
            "新 chunk 目录支持更细粒度的范围查询。",
            "SDSL 尾区里的 zero/copy 位图帮助快速判定 GT 行类型。",
            "copy 行还要结合 copy-origin map 恢复 unique 行。",
            "lossless 样本查询路径可只输出 fixed fields + GT，不强制还原 other fields。",
        ],
        COLORS["soft_blue"],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "5. 汇报建议", "风险点与下一步")
    add_bullet_box(
        slide,
        Inches(0.75),
        Inches(1.2),
        Inches(4.0),
        Inches(4.85),
        "风险点",
        [
            "顶层没有全局 `magic/version`，格式自描述能力不足。",
            "archive 中含 `bool`、`size_t`，跨平台 ABI 稳定性不足。",
            "final `.gsc` 中 chunk 依赖外层索引定位，不是完全自包含。",
        ],
        COLORS["soft_orange"],
        title_color=COLORS["orange"],
    )
    add_bullet_box(
        slide,
        Inches(4.95),
        Inches(1.2),
        Inches(3.9),
        Inches(4.85),
        "建议先做",
        [
            "拿一个 toy 文件做偏移级 hexdump。",
            "补一轮 legacy chunk 老文件实测。",
            "把写路径与读路径锚点整理进正式说明文档。",
        ],
        COLORS["soft_teal"],
        title_color=COLORS["teal"],
    )
    add_bullet_box(
        slide,
        Inches(9.0),
        Inches(1.2),
        Inches(3.55),
        Inches(4.85),
        "若走协议化",
        [
            "先引入顶层 `magic/version`。",
            "固定字段宽度，去掉 ABI 敏感类型。",
            "再考虑与 AVS 外层容器做正式映射。",
        ],
        COLORS["soft_blue"],
    )
    add_footer(slide)

    prs.save(path)


def build_avs_deck(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title_slide(
        slide,
        "VCF / gVCF 接入\nAVS-G 映射汇报",
        "基于 VCF_gVCF接入AVS映射汇报-260416.html\n聚焦接入思路、边界与落地顺序",
        "简洁版 / 6 页",
    )

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "1. 总体判断", "方向对、边界清、可分阶段落地")
    add_bullet_box(
        slide,
        Inches(0.75),
        Inches(1.2),
        Inches(6.05),
        Inches(4.95),
        "这份映射草案为什么值得做",
        [
            "把 AVS 外层容器和内部 codec 分开讨论，避免一上来重写算法。",
            "VCF 先按 chunk 对齐 AVS block，改动小且能复用现有索引。",
            "gVCF 先保 native V4 block codec，只替换外层封装，迁移风险最低。",
            "把 header / sample / ploidy / n_samples 上移到 HEADER，语义更干净。",
        ],
        COLORS["soft_blue"],
    )
    add_callout(
        slide,
        Inches(7.1),
        Inches(1.35),
        Inches(5.35),
        Inches(1.15),
        "三层关系",
        "AVS-G 是封装框架，现有 `.gsc / gVCF` 是业务 payload，映射草案是两者之间的桥。",
        COLORS["teal"],
    )
    add_callout(
        slide,
        Inches(7.1),
        Inches(2.85),
        Inches(5.35),
        Inches(1.15),
        "草案的核心价值",
        "把块数据、全局索引与参数重新归位，为分阶段 reader / writer 改造提供施工图。",
        COLORS["orange"],
    )
    add_callout(
        slide,
        Inches(7.1),
        Inches(4.35),
        Inches(5.35),
        Inches(1.15),
        "重点不是换算法",
        "重点是把现有 payload 改写成更稳定、更显式的协议表达。",
        COLORS["navy"],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "2. VCF 映射主线", "这部分已经相对成型")
    stages = [
        {"title": "HEADER", "body": "STD_TYPE\nLOSS_MODE\nHEADER_TEXT\nSAMPLE_LIST", "fill": COLORS["navy"]},
        {"title": "COMP_DATA", "body": "chunk 对齐\nBLOCK_FIXED_FIELDS\nBLOCK_GT", "fill": COLORS["teal"]},
        {"title": "TAILER", "body": "BLOCK_INFO\nPARAM_GT\nPARAM_OTHER_FIELDS", "fill": COLORS["orange"]},
    ]
    add_stage_row(slide, Inches(1.55), stages)
    add_bullet_box(
        slide,
        Inches(0.8),
        Inches(3.15),
        Inches(5.8),
        Inches(2.95),
        "VCF 映射重点",
        [
            "首版把 `VCF_BLOCK_INDIV` 固定到 chunk 粒度，不直接下钻 row_block。",
            "fixed fields、GT 放入 block payload；lossless 时可先接受 whole-blob `BLOCK_OTHERFIELDS`。",
            "chunk 偏移、位置、染色体边界等全局索引迁入 `BLOCK_INFO`。",
        ],
        COLORS["soft_teal"],
        title_color=COLORS["teal"],
    )
    add_bullet_box(
        slide,
        Inches(6.95),
        Inches(3.15),
        Inches(5.55),
        Inches(2.95),
        "为什么这样映射更稳",
        [
            "最大化复用当前 `.gsc` 的 chunk 组织与 reader 认知。",
            "把实现细节从裸 offset/bool 中剥离，改成显式协议节点。",
            "先把外层包通，再逐步细化块内结构，工程风险更低。",
        ],
        COLORS["soft_blue"],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "3. gVCF 映射主线", "先外包，不强拆")
    add_bullet_box(
        slide,
        Inches(0.75),
        Inches(1.2),
        Inches(5.95),
        Inches(5.0),
        "为什么保 native block 是对的",
        [
            "现有 gVCF 模块已经有独立 magic / version / block index 体系。",
            "它最关键的能力是范围查询，不是为了追求节点对称去拆碎内部 codec。",
            "短期应优先保住 native block + index，自上而下接入 AVS 外层。",
        ],
        COLORS["soft_orange"],
        title_color=COLORS["orange"],
    )
    add_bullet_box(
        slide,
        Inches(6.95),
        Inches(1.2),
        Inches(5.55),
        Inches(5.0),
        "建议映射",
        [
            "`BLOCK_CORE_FIELDS`：保留 native block codec。",
            "`BLOCK_SAMPLE_FIELDS`：承接现有字段 hybrid 编码。",
            "`BLOCK_INDEX`：保留块级范围查询索引。",
            "`PARAM_GVCF`：保存内部版本号与块配置。",
        ],
        COLORS["soft_blue"],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "4. 关键未闭环问题", "已经从概念问题收敛到工程设计问题")
    add_bullet_box(
        slide,
        Inches(0.75),
        Inches(1.2),
        Inches(3.9),
        Inches(5.0),
        "问题 1-2",
        [
            "`BLOCK_OTHERFIELDS` 首版是 whole-blob 还是块内化。",
            "GT zero/copy 与 copy-origin map 如何写成 `PARAM_GT`。",
        ],
        COLORS["soft_orange"],
        title_color=COLORS["orange"],
    )
    add_bullet_box(
        slide,
        Inches(4.85),
        Inches(1.2),
        Inches(3.9),
        Inches(5.0),
        "问题 3-4",
        [
            "HEADER_TEXT、BLOCK_GT、OTHERFIELDS 是否共享统一 backend registry。",
            "当前 `.gsc` 里的 `bool`、`size_t`、匿名尾区不能直接升级成对外协议字段。",
        ],
        COLORS["soft_teal"],
        title_color=COLORS["teal"],
    )
    add_bullet_box(
        slide,
        Inches(8.95),
        Inches(1.2),
        Inches(3.6),
        Inches(5.0),
        "问题 5-6",
        [
            "现有 `other_fields_offset / sdsl_offset` 这套 offset 观念应由 AVS 外层 BLOCK_INFO 接管。",
            "首版必须把块粒度固定到 chunk，否则 reader / writer 接口会持续漂移。",
        ],
        COLORS["soft_blue"],
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    add_full_bg(slide, COLORS["white"])
    add_top_band(slide, "5. 建议落地顺序", "先包外层，再迁 payload，再迁索引参数")
    stages = [
        {"title": "Phase 1", "body": "迁 HEADER\n类型、模式、样本与 ploidy", "fill": COLORS["navy"]},
        {"title": "Phase 2", "body": "按 chunk 封装\nfixed fields / GT / otherfields", "fill": COLORS["teal"]},
        {"title": "Phase 3", "body": "迁 TAILER\nBLOCK_INFO 与参数区", "fill": COLORS["orange"]},
        {"title": "Phase 4", "body": "再优化\n块内化与统一 backend", "fill": COLORS["gold"]},
    ]
    add_stage_row(slide, Inches(1.75), stages)
    add_callout(
        slide,
        Inches(0.95),
        Inches(4.0),
        Inches(11.4),
        Inches(1.55),
        "为什么这个顺序更稳",
        "先把 AVS 外层语义包起来，再逐步迁移 payload 和索引参数，最后再做块内细化与统一 registry。这样能尽快得到可用 reader / writer 原型，同时避免在一开始就碰最难的协议细节。",
        COLORS["teal"],
    )
    add_footer(slide)

    prs.save(path)


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_gsc_deck(OUT_DIR / "GSC多样本二进制格式汇报-简洁版-260416.pptx")
    build_avs_deck(OUT_DIR / "VCF_gVCF接入AVS映射汇报-简洁版-260416.pptx")


if __name__ == "__main__":
    main()
